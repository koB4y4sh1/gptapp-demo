from pptx.presentation import Presentation
from pptx.util import Inches
from src.domain.model.slides.base_slide import BaseSlide
import os

class ImageSlide(BaseSlide):
    def __init__(self, prs: Presentation, header: str, content: str, image_path: str = None):
        super().__init__(prs, header, content)
        self.image_path = image_path

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content
        
        if not self.image_path:
            return

        if  os.path.exists(self.image_path):
            # ローカルファイルの場合はそのまま挿入
            slide.shapes.add_picture(self.image_path, Inches(1), Inches(1.5), width=Inches(4))
        else:
            # 無効なパスの場合はメッセージを追加
            slide.placeholders[1].text += f"\n[画像パス無効: {self.image_path}]"