from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.presentation import Presentation
from typing import List
from src.domain.model.slides.base_slide import BaseSlide

class ThreeHorizontalFlowSlide(BaseSlide):
    def __init__(self, prs: Presentation, header: str, content: str, steps: List[str]):
        super().__init__(prs, header, content)
        self.steps = steps

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content

        # ステップを横並びで配置
        for i, step in enumerate(self.steps):
            left = Inches(1 + i * 3)
            top = Inches(1.5)
            width = Inches(2)
            height = Inches(1)
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.text = step
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(91, 155, 213)  # 青色
            shape.shadow.inherit = False
            shape.shadow.blur_radius = Inches(0.1)
            shape.shadow.distance = Inches(0.1)
