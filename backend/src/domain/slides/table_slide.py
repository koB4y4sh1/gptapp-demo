from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor
from src.domain.slides.base import Base

class TableSlide(Base):
    def __init__(self, prs, header, content, table_data):
        super().__init__(prs, header, content)
        self.table_data = table_data

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content

        # テーブルの行数と列数を取得
        rows = len(self.table_data)
        cols = len(self.table_data[0])

        # テーブルを追加
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table

        # セルにデータを入力
        for row_idx, row in enumerate(self.table_data):
            for col_idx, cell_data in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                # セルのフォント設定
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
