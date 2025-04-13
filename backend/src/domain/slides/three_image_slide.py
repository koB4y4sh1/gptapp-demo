from pptx.util import Inches
from src.domain.slides.base import Base

class ThreeImageSlide(Base):
    def __init__(self, prs, header, content, image_paths):
        super().__init__(prs, header, content)
        self.image_paths = image_paths

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content

        # 画像を横並びで配置
        for i, image_path in enumerate(self.image_paths):
            slide.shapes.add_picture(image_path, Inches(1 + i * 3), Inches(1.5), width=Inches(2))
