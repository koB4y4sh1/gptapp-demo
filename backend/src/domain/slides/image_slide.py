from pptx import Presentation
from pptx.util import Inches
from backend.src.domain.slides.base import Base

class ImageSlide(Base):
    def __init__(self, prs: Presentation, header: str, content: str, image_path: str = None):
        super().__init__(prs, header, content)
        self.image_path = image_path

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content
        if self.image_path:
            slide.shapes.add_picture(self.image_path, Inches(1), Inches(1.5), width=Inches(4))