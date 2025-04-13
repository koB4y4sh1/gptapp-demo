from pptx import Presentation
from src.domain.slides.base import Base

class TextSlide(Base):
    def __init__(self, prs: Presentation, header: str, content: str):
        super().__init__(prs, header, content)

    def build(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])
        slide.shapes.title.text = self.header
        slide.placeholders[1].text = self.content