from pptx import Presentation
from io import BytesIO
from typing import List
from src.domain.model.slides.text_slide import TextSlide
from src.domain.model.slides.image_slide import ImageSlide
from src.domain.model.slides.three_image_slide import ThreeImageSlide
from src.domain.model.slides.table_slide import TableSlide
from src.domain.model.slides.three_horizontal_flow_slide import ThreeHorizontalFlowSlide
from src.domain.model.type.page import Page
from src.domain.model.type.template import TemplateType
from src.utils.logger import get_logger

logger = get_logger(__name__)

def generate_pptx(pages: List[Page]) -> bytes:
    """
    ページリストからPPTXバイナリを生成

    Args:
        pages (List[Page]): スライドページ情報リスト

    Returns:
        bytes: PPTXファイルのバイナリ
    """
    prs = Presentation()

    for page in pages:
        if page.template == TemplateType.TEXT:
            slide = TextSlide(prs, page.header, page.content)
        elif page.template == TemplateType.IMAGE:
            slide = ImageSlide(prs, page.header, page.content, page.images[0] if page.images else None)
        elif page.template == TemplateType.THREE_IMAGES:
            slide = ThreeImageSlide(prs, page.header, page.content, page.images)
        elif page.template == TemplateType.TABLE:
            slide = TableSlide(prs, page.header, page.content, page.table)
        elif page.template == TemplateType.THREE_HORIZONTAL_FLOW:
            slide = ThreeHorizontalFlowSlide(prs, page.header, page.content, page.steps)
        else:
            logger.warning(f"未対応テンプレート: {page.template}")
            continue

        slide.build()

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()

if __name__ == "__main__":
    # テスト用ダミーデータ作成
    pages = [
        Page(
            header="Pythonとは",
            content="Pythonの概要を説明する",
            template=TemplateType.TEXT,
            images=None,
            captions=None,
            table=None,
            steps=None
        ),
        Page(
            header="プログラムの歴史",
            content="Pythonの歴史と発展を紹介する",
            template=TemplateType.IMAGE,
            images=[],
            captions=["説明するイラスト"],
            table=None,
            steps=None
        ),
        Page(
            header="活用分野",
            content="PythonはWeb開発やデータ分析など様々な分野で使われています。",
            template=TemplateType.THREE_IMAGES,
            images=[],
            captions=["Pythonについてのイラスト", "スマートデバイスを操作するイラスト", "説明するイラスト"],
            table=None,
            steps=None
        ),
        Page(
            header="他言語との比較",
            content="以下は主要な言語との比較表です。",
            template=TemplateType.TABLE,
            images=[],
            captions=[],
            table=[
                ["言語", "用途", "学習難易度"],
                ["Python", "汎用", "易しい"],
                ["Java", "業務アプリ", "中"],
                ["C++", "システム", "難しい"]
            ],
            steps=None
        ),
    ]

    result = generate_pptx(pages)
    print(f"実行結果 (slide_state): {result}")
