from pptx import Presentation
from io import BytesIO
from src.domain.slides.text_slide import TextSlide
from src.domain.slides.image_slide import ImageSlide
from src.domain.slides.three_image_slide import ThreeImageSlide
from src.domain.slides.table_slide import TableSlide
from src.domain.slides.three_horizontal_flow_slide import ThreeHorizontalFlowSlide

def generate_ppt(data: dict) -> bytes:
    prs = Presentation()
    slides = data.get("pages", [])

    for page in slides:
        header = page.get("header", "")
        content = page.get("content", "")
        template = page.get("template", "text")
        image_paths = page.get("images", [])
        table_data = page.get("table", [])
        steps = page.get("steps", [])

        if template == "text":
            slide = TextSlide(prs, header, content)
        elif template == "image":
            slide = ImageSlide(prs, header, content, image_paths[0] if image_paths else None)
        elif template == "three_images":
            slide = ThreeImageSlide(prs, header, content, image_paths)
        elif template == "table":
            slide = TableSlide(prs, header, content, table_data)
        elif template == "three_horizontal_flow":
            slide = ThreeHorizontalFlowSlide(prs, header, content, steps)
        else:
            continue

        slide.build()

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
