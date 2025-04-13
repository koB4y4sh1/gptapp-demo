# generate.py

from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def generate_ppt(data: dict) -> bytes:
    """
    入力データ（JSON）をもとにPowerPointファイルを生成してバイト列で返す。
    """
    prs = Presentation()

    slides:list[dict] = data.get("slides",[])
    for sile_data  in slides:
        section_title_text = sile_data.get("section", "")
        content_text = sile_data.get("content", "")
        image_path = sile_data.get("image")  # 任意。画像を挿入したい場合。

        slide = prs.slides.add_slide(prs.slide_layouts[1])  # タイトルとコンテンツレイアウト
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = section_title_text
        content.text = content_text

        if image_path:
            try:
                # 左2インチ, 上3インチの位置に、幅4インチで画像を配置
                slide.shapes.add_picture(image_path, Inches(2), Inches(3), width=Inches(4))
            except Exception as e:
                print(f"画像挿入エラー: {e}")

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()
